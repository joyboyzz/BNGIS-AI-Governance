"""
BNGIS — College presentation builder (16:9, dark theme, tricolor accents)
Output: docs/BNGIS_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- palette ----------------
BG      = RGBColor(0x0B, 0x10, 0x20)   # deep navy
BG2     = RGBColor(0x10, 0x18, 0x38)
CARD    = RGBColor(0x16, 0x20, 0x4A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x97, 0xA0, 0xBD)
SAFFRON = RGBColor(0xFF, 0x99, 0x33)
GREEN   = RGBColor(0x2E, 0xB8, 0x72)
INDIGO  = RGBColor(0x7C, 0x8C, 0xFF)
CYAN    = RGBColor(0x38, 0xD6, 0xD0)
PINK    = RGBColor(0xFF, 0x6B, 0x9D)
RED     = RGBColor(0xFF, 0x5C, 0x5C)
GOLD    = RGBColor(0xFF, 0xD6, 0x66)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=CARD, line=None, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try: shp.adjustments[0] = 0.08
        except Exception: pass
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1.2)
    else:
        shp.line.fill.background()
    return shp


def text(s, x, y, w, h, runs, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.0):
    """runs: str OR list of (txt, {size, color, bold}) paragraph tuples OR
    list of lists for multi-run paragraphs."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for para in runs:
        if isinstance(para, tuple):
            para = [para]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = leading
        for txt, st in para:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(st.get("size", size))
            f.color.rgb = st.get("color", color)
            f.bold = st.get("bold", bold)
    return tb


def kicker(s, label, color=SAFFRON):
    text(s, 0.7, 0.42, 10, 0.5, [(label, {"size": 15, "color": color,
        "bold": True})])

def title(s, t, color=WHITE, size=33):
    text(s, 0.7, 0.78, 12.2, 0.95, [(t, {"size": size, "color": color,
        "bold": True})])

def tricolor_bar(s, y=7.28):
    box(s, 0,    y, W / 3,     0.22, SAFFRON, radius=False)
    box(s, W / 3, y, W / 3,    0.22, RGBColor(0xF5, 0xF5, 0xF5), radius=False)
    box(s, 2 * W / 3, y, W / 3, 0.22, RGBColor(0x13, 0x88, 0x08), radius=False)

def pageno(s, n):
    text(s, 12.35, 6.98, 0.8, 0.4, [(str(n), {"size": 12, "color": MUTED})],
         align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, h, label, fill=CARD, fg=WHITE, size=14, bold=True,
         border=None):
    box(s, x, y, w, h, fill=fill, line=border)
    text(s, x, y, w, h, [(label, {"size": size, "color": fg, "bold": bold})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def bullet_list(s, x, y, w, items, size=16, gap=0.52, color=WHITE,
                marker="✦ ", mcolor=SAFFRON):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            it, st = it
        else:
            st = {}
        text(s, x, y + i * gap, w, gap,
             [(marker, {"size": size, "color": mcolor, "bold": True}),
              (it, {"size": st.get("size", size),
                    "color": st.get("color", color),
                    "bold": st.get("bold", False)})])


# ==========================================================================
# SLIDE 1 — TITLE
# ==========================================================================
s = slide()
# glow circles
box(s, 9.2, -1.5, 6, 6, fill=RGBColor(0x14, 0x1D, 0x42))
box(s, 10.2, -0.5, 4, 4, fill=RGBColor(0x18, 0x24, 0x52))
# neural dots (right side)
import random
rng = random.Random(11)
for _ in range(26):
    x = rng.uniform(9.0, 13.0); y = rng.uniform(0.6, 6.6)
    d = rng.uniform(0.06, 0.16)
    c = rng.choice([SAFFRON, INDIGO, CYAN, GREEN])
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                           Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = c; o.line.fill.background()
    o.shadow.inherit = False
text(s, 0.8, 0.85, 9, 0.5, [("BHARATH  •  NEURO  •  GOVERNANCE  •  INTELLIGENCE",
    {"size": 15, "color": INDIGO, "bold": True})])
text(s, 0.75, 1.35, 9.5, 1.9, [
    ("BNGIS", {"size": 88, "color": SAFFRON, "bold": True})])
text(s, 0.8, 3.05, 10.5, 1.0, [
    ("An AI brain that connects every citizen to every government "
     "resource — in real time.", {"size": 22, "color": WHITE})])
chip(s, 0.8,  4.35, 2.15, 0.52, "8 AI Modules", border=INDIGO, fg=RGBColor(0xB9,0xC3,0xFF))
chip(s, 3.15, 4.35, 1.95, 0.52, "25 Tests ✓",  border=GREEN,  fg=RGBColor(0x9F,0xE8,0xC2))
chip(s, 5.3,  4.35, 2.05, 0.52, "4 Languages", border=CYAN,    fg=RGBColor(0x9F,0xE8,0xE4))
chip(s, 7.5,  4.35, 1.75, 0.52, "₹0 Cost",     border=SAFFRON, fg=RGBColor(0xFF,0xC8,0x99))
box(s, 0.8, 5.45, 5.6, 0.02, fill=RGBColor(0x22, 0x30, 0x5E))
text(s, 0.8, 5.65, 11, 1.2, [
    ("Final Year Project  •  2026", {"size": 15, "color": MUTED}),
    ("Praveen Sudireddy  (joyboyzz)", {"size": 18, "color": WHITE, "bold": True}),
    ("github.com/joyboyzz/BNGIS-AI-Governance", {"size": 14, "color": CYAN}),
], leading=1.25)
tricolor_bar(s)

# ==========================================================================
# SLIDE 2 — THE PROBLEM
# ==========================================================================
s = slide()
kicker(s, "THE PROBLEM")
title(s, "India's governance crisis, in numbers")
stats = [
    ("₹3.8 L Cr", "schemes undelivered\nevery year", RED),
    ("67%", "of eligible citizens unaware\nof their schemes", SAFFRON),
    ("₹1.2 L Cr", "leaked to corruption\nin delivery", RED),
    ("23%", "ghost beneficiaries in\nration system", GOLD),
    ("34%", "hospital beds sit\nempty daily", SAFFRON),
    ("847 days", "avg time to fix\none pothole", RED),
]
for i, (v, l, c) in enumerate(stats):
    x = 0.7 + (i % 3) * 4.1
    y = 2.05 + (i // 3) * 2.45
    box(s, x, y, 3.75, 2.1)
    text(s, x + 0.25, y + 0.28, 3.3, 0.9,
         [(v, {"size": 34, "color": c, "bold": True})])
    text(s, x + 0.25, y + 1.15, 3.3, 0.9,
         [(l, {"size": 14.5, "color": MUTED})], leading=1.1)
text(s, 0.7, 6.75, 12, 0.5, [
    ("Source: BNGIS project specification — aggregated from govt data / NDMA / NITI Aayog reports",
     {"size": 11.5, "color": RGBColor(0x5D, 0x66, 0x84)})])
tricolor_bar(s); pageno(s, 2)

# ==========================================================================
# SLIDE 3 — WHY EXISTING SYSTEMS FAIL
# ==========================================================================
s = slide()
kicker(s, "THE GAP")
title(s, "Existing platforms give access — not intelligence")
rows = [
    ("MyGov / Grievance portals", "Feedback only — no intelligence, no prediction"),
    ("UMANG app", "Service access — no optimization or matching"),
    ("DigiLocker", "Documents — no resource mapping"),
    ("Smart City projects", "Siloed — no cross-city / cross-domain learning"),
]
y = 2.0
for name, why in rows:
    box(s, 0.7, y, 7.3, 0.92)
    text(s, 0.95, y, 3.6, 0.92, [(name, {"size": 15.5, "bold": True})],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.6, y, 3.3, 0.92, [(why, {"size": 13, "color": MUTED})],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 1.1
box(s, 8.35, 2.0, 4.25, 4.32, fill=RGBColor(0x1A, 0x14, 0x30), line=PINK)
text(s, 8.65, 2.3, 3.7, 3.8, [
    ("THE MISSING LAYER", {"size": 13, "color": PINK, "bold": True}),
    ("", {}),
    ("None of them create an intelligent neural network that connects "
     "citizens → schemes → resources → audits…",
     {"size": 15, "color": WHITE}),
    ("", {}),
    ("…AUTOMATICALLY.", {"size": 19, "color": SAFFRON, "bold": True}),
], leading=1.15)
tricolor_bar(s); pageno(s, 3)

# ==========================================================================
# SLIDE 4 — THE SOLUTION / MODULE MAP
# ==========================================================================
s = slide()
kicker(s, "THE SOLUTION", GREEN)
title(s, "BNGIS — one platform, eight intelligence modules", size=31)
mods = [
    ("1", "Citizen Neural Profile", "10-dim need vector, privacy-first", INDIGO, "✅"),
    ("2", "Resource Optimization", "allocation + reroute advisories", CYAN, "✅"),
    ("3", "Scheme Matching (SMDE)", "21 real schemes, knapsack portfolio", SAFFRON, "✅"),
    ("4", "Corruption Shield (CDS)", "Benford + 5-layer anomaly AI", RED, "✅"),
    ("5", "Disaster Response (DRNN)", "SEIR simulator + flood risk", CYAN, "✅"),
    ("7", "Transparency Blockchain", "SHA-256 chain, tamper-proof", GREEN, "✅"),
    ("8", "Citizen Voice NLP", "4-language grievance AI", PINK, "✅"),
    ("10", "Public Dashboard", "bilingual, real-time", INDIGO, "✅"),
]
for i, (num, name, desc, c, tick) in enumerate(mods):
    x = 0.7 + (i % 4) * 3.12
    y = 2.1 + (i // 4) * 2.3
    box(s, x, y, 2.9, 2.0)
    chip(s, x + 0.22, y + 0.22, 0.62, 0.62, num, fill=c, fg=BG, size=20)
    text(s, x + 0.22, y + 0.95, 2.5, 0.5,
         [(name, {"size": 14.5, "bold": True})])
    text(s, x + 0.22, y + 1.38, 2.5, 0.6,
         [(desc, {"size": 11.5, "color": MUTED})], leading=1.05)
text(s, 0.7, 6.7, 12, 0.5, [
    ("✅ = implemented & live in this MVP   ·   Modules 6 & 9 (predictive governance, "
     "coordination) on the roadmap", {"size": 13, "color": MUTED})])
tricolor_bar(s); pageno(s, 4)

# ==========================================================================
# SLIDE 5 — ARCHITECTURE
# ==========================================================================
s = slide()
kicker(s, "ARCHITECTURE", INDIGO)
title(s, "Zero-dependency intelligence stack", size=31)
# browser box
box(s, 0.7, 1.95, 5.4, 1.15, line=INDIGO)
text(s, 0.9, 2.05, 5.0, 0.95, [
    ("🌐 BROWSER — SPA UI", {"size": 14, "color": INDIGO, "bold": True}),
    ("Vanilla JS + SVG charts · EN/తెలుగు · 8 pages · no CDN",
     {"size": 12, "color": MUTED})], leading=1.15)
# api box
box(s, 0.7, 3.35, 5.4, 1.15, line=SAFFRON)
text(s, 0.9, 3.45, 5.0, 0.95, [
    ("⚡ FASTAPI REST LAYER", {"size": 14, "color": SAFFRON, "bold": True}),
    ("12 endpoints · auto docs · CORS · pydantic validation",
     {"size": 12, "color": MUTED})], leading=1.15)
# engines box
box(s, 0.7, 4.75, 5.4, 1.95, line=GREEN)
text(s, 0.9, 4.85, 5.0, 1.8, [
    ("🧠 AI ENGINE CORE (pure Python)", {"size": 14, "color": GREEN, "bold": True}),
    ("cosine matching · knapsack DP · Benford χ² · robust z-scores · HHI ·",
     {"size": 12, "color": MUTED}),
    ("SEIR ODEs · haversine · fuzzy SequenceMatcher · SHA-256 chaining",
     {"size": 12, "color": MUTED})], leading=1.25)
# arrows
for yy in (3.1, 4.5, 5.7):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.2), Inches(yy),
                           Inches(0.35), Inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background()
    a.shadow.inherit = False
# right: data + quality
box(s, 6.6, 1.95, 6.0, 1.5, fill=RGBColor(0x11, 0x2A, 0x22), line=GREEN)
text(s, 6.85, 2.08, 5.6, 1.3, [
    ("🗂️ REAL DATA LAYER", {"size": 14, "color": GREEN, "bold": True}),
    ("21 real govt schemes (Central + Karnataka) · 24 Mysuru facilities ·",
     {"size": 12, "color": MUTED}),
    ("12 Karnataka districts · 1,250+ seeded audit transactions",
     {"size": 12, "color": MUTED})], leading=1.2)
box(s, 6.6, 3.65, 6.0, 1.5, fill=RGBColor(0x241, 0x1A, 0x2E) if False else RGBColor(0x2A, 0x14, 0x28), line=PINK)
text(s, 6.85, 3.78, 5.6, 1.3, [
    ("⛓️ TRANSPARENCY CHAIN", {"size": 14, "color": PINK, "bold": True}),
    ("every AI action auto-recorded as an immutable block ·",
     {"size": 12, "color": MUTED}),
    ("tamper simulation + one-click repair",
     {"size": 12, "color": MUTED})], leading=1.2)
box(s, 6.6, 5.35, 6.0, 1.35, line=GOLD)
text(s, 6.85, 5.48, 5.6, 1.1, [
    ("🧪 QUALITY GATE", {"size": 14, "color": GOLD, "bold": True}),
    ("25 pytest engine tests · GitHub Actions CI (3.11/3.12/3.13) · Docker",
     {"size": 12, "color": MUTED})], leading=1.2)
tricolor_bar(s); pageno(s, 5)

# ==========================================================================
# SLIDE 6 — SCHEME MATCHING ENGINE
# ==========================================================================
s = slide()
kicker(s, "MODULE 3 — SCHEME MATCHING & DELIVERY ENGINE")
title(s, "BNGIS-Match: the algorithm that finds every rupee", size=30)
steps = [
    ("1 · VECTORIZE", "citizen → 10-dim need vector\nincome · area · occupation · gender · age · caste · disability"),
    ("2 · SIMILARITY", "cosine similarity vs every\nscheme's eligibility signature"),
    ("3 · HARD RULES", "age / income / land / state /\ncaste / flags → pass or blocker"),
    ("4 · PRIORITY", "0.35·benefit + 0.25·success\n+ 0.20·effort + 0.10·sim + 0.10·urgency"),
    ("5 · PORTFOLIO", "greedy knapsack, effort budget,\nconflict rules (MUDRA ✕ Stand-Up)"),
]
x = 0.7
for i, (hd, bd) in enumerate(steps):
    box(s, x, 2.0, 2.35, 2.5)
    text(s, x + 0.15, 2.15, 2.05, 0.5, [(hd, {"size": 12.5, "color": SAFFRON, "bold": True})])
    text(s, x + 0.15, 2.62, 2.05, 1.8, [(bd, {"size": 11, "color": MUTED})], leading=1.12)
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.38),
                               Inches(3.05), Inches(0.28), Inches(0.24))
        a.fill.solid(); a.fill.fore_color.rgb = INDIGO; a.line.fill.background()
        a.shadow.inherit = False
    x += 2.62
# results strip
box(s, 0.7, 4.85, 12.0, 1.7, fill=RGBColor(0x0F, 0x2A, 0x1C), line=GREEN)
text(s, 1.0, 5.0, 11.5, 1.45, [
    ("LIVE DEMO RESULT — “Lakshmi, 34” (rural Karnataka, SC, daily-wage, ₹96k/yr)",
     {"size": 15, "color": GREEN, "bold": True}),
    ("9 eligible schemes found in 90 ms  ·  Optimal portfolio of 4: PM-JAY (₹5L health) + "
     "PMAY-G (pucca house) + Gruha Lakshmi (₹24k/yr) + Shakti (free bus)",
     {"size": 13.5, "color": WHITE}),
    ("≈ ₹44,000+/year in measurable value · each match explains WHY — and why 12 others were blocked",
     {"size": 13.5, "color": MUTED})], leading=1.3)
tricolor_bar(s); pageno(s, 6)

# ==========================================================================
# SLIDE 7 — RESOURCE OPTIMIZATION
# ==========================================================================
s = slide()
kicker(s, "MODULE 2 — RESOURCE OPTIMIZATION CORTEX", CYAN)
title(s, "Ending the “full hospital next to empty beds” paradox", size=30)
box(s, 0.7, 2.0, 6.1, 2.3)
text(s, 0.95, 2.15, 5.6, 2.0, [
    ("ALLOCATION SCORE", {"size": 13.5, "color": CYAN, "bold": True}),
    ("score = 0.55·proximity + 0.30·availability + 0.15·quality",
     {"size": 16, "color": WHITE, "bold": True}),
    ("· haversine distance from citizen location\n· live utilization of every facility\n"
     "· quality rating weighted in", {"size": 13, "color": MUTED})], leading=1.3)
box(s, 0.7, 4.55, 6.1, 2.0, fill=RGBColor(0x2A, 0x18, 0x10), line=SAFFRON)
text(s, 0.95, 4.7, 5.6, 1.75, [
    ("REROUTE ADVISORY (auto)", {"size": 13.5, "color": SAFFRON, "bold": True}),
    ("If nearest facility > 85% full →\nBNGIS redirects to the best facility with "
     "spare capacity + wait-time savings.", {"size": 13.5, "color": WHITE})], leading=1.25)
# right: live numbers
box(s, 7.1, 2.0, 5.55, 4.55, fill=RGBColor(0x0F, 0x1C, 0x2A), line=INDIGO)
text(s, 7.35, 2.15, 5.0, 4.2, [
    ("MYSURU NETWORK (demo)", {"size": 13.5, "color": INDIGO, "bold": True}),
    ("", {}),
    ("24 facilities mapped — hospitals, schools, water works",
     {"size": 14, "color": WHITE}),
    ("", {}),
    ("Vijayanagar → nearest hospital 91% full", {"size": 13.5, "color": RED}),
    ("BNGIS reroutes +2.4 km → 32 beds free, wait −40 min",
     {"size": 13.5, "color": GREEN, "bold": True}),
    ("", {}),
    ("1,075 free beds visible network-wide — the system KNOWS where the "
     "empty beds are", {"size": 13.5, "color": MUTED})], leading=1.22)
tricolor_bar(s); pageno(s, 7)

# ==========================================================================
# SLIDE 8 — CORRUPTION SHIELD
# ==========================================================================
s = slide()
kicker(s, "MODULE 4 — CORRUPTION DETECTION SHIELD", RED)
title(s, "Five AI layers that audit 1,250+ transactions in 120 ms", size=30)
layers = [
    ("L1", "Benford's Law", "first-digit χ² test vs log distribution", RED),
    ("L2", "Statistical anomalies", "robust z-scores (median/MAD)", SAFFRON),
    ("L3", "Vendor network", "share, HHI cartel index, pattern flags", GOLD),
    ("L4", "Temporal patterns", "threshold-split, weekend, FYE spike", PINK),
    ("L5", "Ghost detection", "fuzzy names ≥86% + address clusters", INDIGO),
]
y = 1.95
for code, name, desc, c in layers:
    box(s, 0.7, y, 6.4, 0.78)
    chip(s, 0.85, y + 0.12, 0.62, 0.55, code, fill=c, fg=BG, size=14)
    text(s, 1.6, y, 2.3, 0.78, [(name, {"size": 14.5, "bold": True})],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.9, y, 3.1, 0.78, [(desc, {"size": 11.5, "color": MUTED})],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 0.94
# right results
box(s, 7.4, 1.95, 5.25, 4.55, fill=RGBColor(0x2A, 0x0F, 0x14), line=RED)
text(s, 7.65, 2.1, 4.8, 4.25, [
    ("CAUGHT THE INJECTED FRAUD ✅", {"size": 14, "color": RED, "bold": True}),
    ("", {}),
    ("PWD — 45.3 HIGH", {"size": 15, "color": SAFFRON, "bold": True}),
    ("12 payments just below ₹50k approval threshold", {"size": 12.5, "color": MUTED}),
    ("", {}),
    ("Education — 47.5 HIGH", {"size": 15, "color": SAFFRON, "bold": True}),
    ("₹1–2L round payments on Saturdays to one vendor", {"size": 12.5, "color": MUTED}),
    ("", {}),
    ("Water Dept — 41.2 MEDIUM · March fiscal spike (2.7σ)",
     {"size": 13, "color": GOLD, "bold": True}),
    ("21 ghost duplicate pairs found · clean depts stayed LOW",
     {"size": 13, "color": GREEN, "bold": True})], leading=1.18)
tricolor_bar(s); pageno(s, 8)

# ==========================================================================
# SLIDE 9 — DISASTER RESPONSE
# ==========================================================================
s = slide()
kicker(s, "MODULE 5 — DISASTER RESPONSE NEURAL NETWORK", CYAN)
title(s, "From reactive chaos to predictive response", size=31)
box(s, 0.7, 2.0, 6.1, 4.55, fill=RGBColor(0x0F, 0x1C, 0x2A), line=CYAN)
text(s, 0.95, 2.15, 5.6, 4.3, [
    ("🦠 SEIR EPIDEMIC SIMULATOR", {"size": 15, "color": CYAN, "bold": True}),
    ("dS/dt = −βSI/N · dE/dt = βSI/N − σE · dI/dt = σE − γI · dR/dt = γI",
     {"size": 13, "color": WHITE}),
    ("", {}),
    ("4 what-if scenarios (0 / 20 / 40 / 60% contact reduction):",
     {"size": 13.5, "color": MUTED}),
    ("No action → peak 162,637 concurrent cases", {"size": 13.5, "color": RED}),
    ("Full response → peak collapses to 1,516  (−99%)",
     {"size": 13.5, "color": GREEN, "bold": True}),
    ("", {}),
    ("CROSS-MODULE LINK → hospital peak demand 8,132 beds vs 4,640 capacity",
     {"size": 13, "color": SAFFRON, "bold": True}),
    ("→ DEFICIT 3,496 beds: field hospitals + district rerouting plan",
     {"size": 13, "color": WHITE})], leading=1.25)
box(s, 7.1, 2.0, 5.55, 4.55, fill=RGBColor(0x0F, 0x2A, 0x24), line=GREEN)
text(s, 7.35, 2.15, 5.05, 4.3, [
    ("🌊 FLOOD RISK — 12 KARNATAKA DISTRICTS", {"size": 15, "color": GREEN, "bold": True}),
    ("risk = 0.35·rain + 0.30·river + 0.20·terrain + 0.15·history",
     {"size": 12.5, "color": MUTED}),
    ("", {}),
    ("SEVERE: Kodagu 91.6 · Uttara Kannada 87.9 · Dakshina Kannada 81.6",
     {"size": 13.5, "color": RED, "bold": True}),
    ("~3.9 lakh people potentially displaced (48h window)",
     {"size": 13, "color": WHITE}),
    ("", {}),
    ("AUTO-GENERATED RESPONSE PLANS:", {"size": 13, "color": GREEN, "bold": True}),
    ("93 shelters · 93 water tankers · medical teams · boats ·\n"
     "vernacular early-warning SMS — per district", {"size": 13, "color": MUTED})], leading=1.25)
tricolor_bar(s); pageno(s, 9)

# ==========================================================================
# SLIDE 10 — BLOCKCHAIN + VOICE NLP
# ==========================================================================
s = slide()
kicker(s, "MODULES 7 & 8 — TRUST LAYER + CITIZEN VOICE", GREEN)
title(s, "Tamper-proof records + grievances in any language", size=30)
box(s, 0.7, 2.0, 6.1, 4.55, fill=RGBColor(0x14, 0x2A, 0x16), line=GREEN)
text(s, 0.95, 2.15, 5.6, 4.3, [
    ("⛓️ TRANSPARENCY BLOCKCHAIN", {"size": 15, "color": GREEN, "bold": True}),
    ("hash = SHA-256(index · time · data · prev_hash · nonce)",
     {"size": 12.5, "color": MUTED}),
    ("", {}),
    ("Every AI action (match / audit / ticket / alert) is appended as an "
     "immutable block — no mining, no crypto, pure integrity.",
     {"size": 13.5, "color": WHITE}),
    ("", {}),
    ("LIVE HACKER SIMULATION:", {"size": 13, "color": RED, "bold": True}),
    ("attacker edits a signed record → verification recomputes every hash → "
     "⛔ TAMPER DETECTED at block N → one-click repair",
     {"size": 13, "color": MUTED})], leading=1.25)
box(s, 7.1, 2.0, 5.55, 4.55, fill=RGBColor(0x2A, 0x10, 0x22), line=PINK)
text(s, 7.35, 2.15, 5.05, 4.3, [
    ("🗣️ CITIZEN VOICE NLP (EN · हिंदी · తెలుగు · ಕನ್ನಡ)", {"size": 14, "color": PINK, "bold": True}),
    ("", {}),
    ("language → sentiment → 9-category intent → urgency → P1–P4 ticket "
     "+ department auto-routing", {"size": 13.5, "color": WHITE}),
    ("", {}),
    ("“राशन कार्ड के लिए दलाल 500 रुपये मांग रहा है”", {"size": 13.5, "color": GOLD, "bold": True}),
    ("→ detected: Hindi · corruption ×4 weighting → routed to LOKAYUKTA, "
     "not the ration office", {"size": 13, "color": MUTED}),
    ("", {}),
    ("Snake-bite emergency → P1 ticket, 1-hour SLA, field dispatch",
     {"size": 13, "color": GREEN, "bold": True})], leading=1.22)
tricolor_bar(s); pageno(s, 10)

# ==========================================================================
# SLIDE 11 — TECH STACK & QUALITY
# ==========================================================================
s = slide()
kicker(s, "ENGINEERING", INDIGO)
title(s, "100% free & open — engineered like production", size=31)
cols = [
    ("BACKEND", INDIGO, [
        "Python 3.11–3.13", "FastAPI + Uvicorn",
        "12 REST endpoints + auto docs", "Zero external AI deps"]),
    ("AI CORE", SAFFRON, [
        "cosine similarity + knapsack", "Benford χ² + z-scores + HHI",
        "SEIR ODE solver", "fuzzy matcher + SHA-256 chain"]),
    ("FRONTEND", CYAN, [
        "vanilla JS SPA — no framework", "hand-rolled SVG charts",
        "English / తెలుగు toggle", "PWA-ready responsive"]),
    ("QUALITY", GREEN, [
        "25 pytest engine tests", "GitHub Actions CI ×3 pythons",
        "Docker + compose", "MIT licensed"]),
]
x = 0.7
for hd, c, items in cols:
    box(s, x, 2.0, 3.0, 3.6)
    chip(s, x + 0.25, 2.25, 2.5, 0.55, hd, fill=c, fg=BG, size=15)
    yy = 3.0
    for it in items:
        text(s, x + 0.25, yy, 2.55, 0.5,
             [("· ", {"size": 12.5, "color": c, "bold": True}),
              (it, {"size": 12.5, "color": MUTED})])
        yy += 0.62
    x += 3.12
box(s, 0.7, 5.85, 12.05, 0.95, fill=RGBColor(0x0F, 0x2A, 0x1C), line=GREEN)
text(s, 1.0, 5.85, 11.5, 0.95, [
    ("TOTAL BUILD & RUN COST: ₹0  —  no paid APIs, no cloud, no licenses. "
     "Deploy: pip install -r requirements.txt && uvicorn app.main:app",
     {"size": 14, "color": GREEN, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)
tricolor_bar(s); pageno(s, 11)

# ==========================================================================
# SLIDE 12 — TESTING
# ==========================================================================
s = slide()
kicker(s, "VALIDATION", GOLD)
title(s, "Tests caught a real epidemiology bug — before GitHub did", size=29)
box(s, 0.7, 2.0, 5.9, 4.55)
text(s, 0.95, 2.2, 5.4, 4.2, [
    ("25/25 PYTEST TESTS — GREEN", {"size": 16, "color": GREEN, "bold": True}),
    ("", {}),
    ("• scheme eligibility ×5 (Lakshmi, farmer, income-block, conflicts)",
     {"size": 13, "color": WHITE}),
    ("• corruption ×5 (determinism, injected fraud, ghosts, Benford)",
     {"size": 13, "color": WHITE}),
    ("• blockchain ×3 (verify, tamper-detect, repair)", {"size": 13, "color": WHITE}),
    ("• resources ×3 (reroute fires, healthy doesn't, ranking order)",
     {"size": 13, "color": WHITE}),
    ("• disaster ×3 (SEIR conservation, −60% effect, Kodagu SEVERE)",
     {"size": 13, "color": WHITE}),
    ("• voice NLP ×6 (4 langs, routing, bribe→Lokayukta, P1, feed)",
     {"size": 13, "color": WHITE})], leading=1.3)
box(s, 6.9, 2.0, 5.75, 2.1, fill=RGBColor(0x2A, 0x18, 0x10), line=SAFFRON)
text(s, 7.15, 2.18, 5.3, 1.8, [
    ("🐛 THE BUG THE TESTS CAUGHT", {"size": 14, "color": SAFFRON, "bold": True}),
    ("SEIR conservation test: S+E+I+R ≠ population — initial E₀ was never "
     "subtracted from S. Silent +240 ghost people at t=0. Fixed, locked with "
     "a regression test.", {"size": 13, "color": MUTED})], leading=1.25)
box(s, 6.9, 4.35, 5.75, 2.2, fill=RGBColor(0x0F, 0x1C, 0x2A), line=INDIGO)
text(s, 7.15, 4.53, 5.3, 1.9, [
    ("☁️ CI — GITHUB ACTIONS", {"size": 14, "color": INDIGO, "bold": True}),
    ("On every push: pytest × {3.11, 3.12, 3.13} + JS syntax check + live "
     "server boot smoke test hitting all 9 endpoint groups.",
     {"size": 13, "color": MUTED}),
    ("Status: ✅ passing (live badge on README)", {"size": 13.5, "color": GREEN, "bold": True})], leading=1.25)
tricolor_bar(s); pageno(s, 12)

# ==========================================================================
# SLIDE 13 — IMPACT
# ==========================================================================
s = slide()
kicker(s, "IMPACT", GREEN)
title(s, "What changes when governance thinks", size=32)
imp = [
    ("₹44k+/yr", "extra benefits per matched\ncitizen family (demo)", GREEN),
    ("−40 min", "wait time via hospital\nreroute advisory", CYAN),
    ("−99%", "epidemic peak with full\nAI-guided response", SAFFRON),
    ("120 ms", "full corruption audit of\n1,250 transactions", PINK),
]
for i, (v, l, c) in enumerate(imp):
    x = 0.7 + i * 3.12
    box(s, x, 2.0, 2.9, 2.05)
    text(s, x + 0.22, 2.2, 2.5, 0.8, [(v, {"size": 30, "color": c, "bold": True})])
    text(s, x + 0.22, 3.05, 2.5, 0.9, [(l, {"size": 12.5, "color": MUTED})], leading=1.1)
box(s, 0.7, 4.4, 12.05, 2.15, fill=RGBColor(0x11, 0x2A, 0x22), line=GREEN)
text(s, 1.0, 4.6, 11.5, 1.8, [
    ("SDG ALIGNMENT", {"size": 13.5, "color": GREEN, "bold": True}),
    ("No Poverty (1) · Zero Hunger (2) · Good Health (3) · Quality Education (4) · "
     "Clean Water (6) · Reduced Inequality (10) · Sustainable Cities (11) · "
     "Climate Action (13) · Peace & Justice (16)",
     {"size": 14.5, "color": WHITE}),
    ("Every rupee saved from leakage is a rupee delivered to the last mile.",
     {"size": 14, "color": SAFFRON, "bold": True})], leading=1.35)
tricolor_bar(s); pageno(s, 13)

# ==========================================================================
# SLIDE 14 — ROADMAP
# ==========================================================================
s = slide()
kicker(s, "WHAT'S NEXT", INDIGO)
title(s, "From Mysuru MVP to national platform", size=32)
phases = [
    ("NOW — MVP", "8 modules live · 1 district · deterministic demo data", GREEN, "✅"),
    ("PHASE 2", "Module 6 Predictive Governance (LSTM + Prophet demand\nforecasting) · Module 9 Coordination Brain", CYAN, "🚧"),
    ("PHASE 3", "PostgreSQL + PostGIS/pgvector · live govt data feeds\n(Airflow ETL from data.gov.in) · auth + RBAC", SAFFRON, "⏳"),
    ("PHASE 4", "State pilot → national scale · Kafka streaming · K3s\n· 12-language voice interface · ₹0 forever", INDIGO, "🌙"),
]
y = 1.98
for hd, bd, c, ic in phases:
    box(s, 0.7, y, 12.0, 1.12)
    chip(s, 0.95, y + 0.26, 1.9, 0.6, hd, fill=c, fg=BG, size=13)
    text(s, 3.1, y, 8.4, 1.12, [(bd, {"size": 13, "color": MUTED})],
         anchor=MSO_ANCHOR.MIDDLE, leading=1.1)
    text(s, 11.9, y, 0.7, 1.12, [(ic, {"size": 20})], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.27
tricolor_bar(s); pageno(s, 14)

# ==========================================================================
# SLIDE 15 — THANK YOU
# ==========================================================================
s = slide(BG2)
for _ in range(30):
    x = rng.uniform(0.3, 13.0); y = rng.uniform(0.4, 7.1)
    d = rng.uniform(0.05, 0.15)
    c = rng.choice([SAFFRON, INDIGO, CYAN, GREEN, PINK])
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = c; o.line.fill.background()
    o.shadow.inherit = False
text(s, 0.8, 2.0, 11.7, 1.2, [("Thank you! 🙏", {"size": 60, "color": SAFFRON, "bold": True})],
     align=PP_ALIGN.CENTER)
text(s, 0.8, 3.3, 11.7, 0.9, [
    ("“Technology should serve democracy, not the other way around.”",
     {"size": 20, "color": WHITE, "italic": True})], align=PP_ALIGN.CENTER)
text(s, 0.8, 4.6, 11.7, 0.6, [
    ("Live demo:  python3 -m uvicorn app.main:app  →  http://localhost:8000",
     {"size": 15, "color": CYAN})], align=PP_ALIGN.CENTER)
text(s, 0.8, 5.3, 11.7, 1.0, [
    ("github.com/joyboyzz/BNGIS-AI-Governance   ·   Built with ❤️ for Bharat",
     {"size": 15, "color": MUTED})], align=PP_ALIGN.CENTER)
text(s, 0.8, 6.3, 11.7, 0.5, [
    ("Praveen Sudireddy (joyboyzz) — Final Year Project 2026",
     {"size": 13, "color": MUTED})], align=PP_ALIGN.CENTER)
tricolor_bar(s)

# ---------------------------------------------------------------------------
out = "/home/user/bngis/docs/BNGIS_Presentation.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
